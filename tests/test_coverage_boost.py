import sys
import unittest
import PIL.Image
from unittest.mock import patch, MagicMock

# Define mocks for third-party libraries that might not be installed locally
sys.modules["openpyxl"] = MagicMock()
sys.modules["pptx"] = MagicMock()
sys.modules["docx"] = MagicMock()
sys.modules["fitz"] = MagicMock()
sys.modules["pytesseract"] = MagicMock()

# Mock PIL Image methods to avoid trying to parse dummy bytes
PIL.Image.open = MagicMock()
PIL.Image.frombytes = MagicMock()


class TestCoverageBoost(unittest.TestCase):
    # Test config
    def test_config(self):
        from config import settings

        self.assertIsNotNone(settings.APP_SECRET_KEY)

    # Test security
    def test_security(self):
        from security import encrypt_data, decrypt_data

        msg = "hello"
        enc = encrypt_data(msg)
        dec = decrypt_data(enc)
        self.assertEqual(dec, msg)

    # Test i18n
    def test_i18n(self):
        from i18n_utils import t, set_language

        res = t("non_existent_key")
        self.assertEqual(res, "non_existent_key")
        set_language("en")
        set_language("invalid_lang")

    # Test local provider
    def test_local_provider(self):
        from local import LocalProvider

        provider = LocalProvider()
        self.assertEqual(provider.name, "Local")
        res = provider.generate_quiz(
            "Python is a programming language. It was created by Guido van Rossum.",
            1,
            question_type="MCQ",
            difficulty="Basic",
            answer_length="1-line",
        )
        self.assertTrue(len(res) > 0)

    # Test gemini provider
    @patch("gemini.generate_ai_quiz")
    def test_gemini_provider(self, mock_gen):
        from gemini import GeminiProvider

        mock_gen.return_value = [{"question": "q"}]
        provider = GeminiProvider(api_key="test")
        self.assertEqual(provider.name, "Gemini")
        res = provider.generate_quiz("text", 1, "MCQ", "Basic", "1-line")
        self.assertEqual(res, [{"question": "q"}])

        with self.assertRaises(ValueError):
            GeminiProvider(api_key="").generate_quiz(
                "text", 1, "MCQ", "Basic", "1-line"
            )

    # Test ollama provider
    @patch("ollama_provider.ollama.Client")
    def test_ollama_provider(self, mock_client_class):
        from ollama_provider import OllamaProvider

        mock_client = MagicMock()
        mock_client.generate.return_value = {
            "response": '{"quiz": [{"question": "q", "type": "MCQ", "answer": "a"}]}'
        }
        mock_client_class.return_value = mock_client

        provider = OllamaProvider(model="llama3", host="http://localhost:11434")
        res = provider.generate_quiz("text", 1, "MCQ", "Basic", "1-line")
        self.assertEqual(len(res), 1)

        # Test error handling
        mock_client.generate.side_effect = Exception("error")
        with self.assertRaises(RuntimeError):
            provider.generate_quiz("text", 1, "MCQ", "Basic", "1-line")

    # Test ollama provider error paths
    @patch("ollama_provider.ollama.Client")
    def test_ollama_provider_errors(self, mock_client_class):
        from ollama_provider import OllamaProvider
        import ollama

        mock_client = MagicMock()
        mock_client_class.return_value = mock_client
        provider = OllamaProvider(model="llama3", host="http://localhost:11434")

        # 1. Empty response
        mock_client.generate.return_value = {"response": ""}
        with self.assertRaises(RuntimeError):
            provider.generate_quiz("text", 1, "MCQ", "Basic", "1-line")

        # 2. JSONDecodeError
        mock_client.generate.return_value = {"response": "invalid-json"}
        with self.assertRaises(RuntimeError):
            provider.generate_quiz("text", 1, "MCQ", "Basic", "1-line")

        # 3. ResponseError model not found
        err = ollama.ResponseError("model not found")
        mock_client.generate.side_effect = err
        with self.assertRaises(RuntimeError):
            provider.generate_quiz("text", 1, "MCQ", "Basic", "1-line")

        # 4. ResponseError other
        err2 = ollama.ResponseError("some other error")
        mock_client.generate.side_effect = err2
        with self.assertRaises(RuntimeError):
            provider.generate_quiz("text", 1, "MCQ", "Basic", "1-line")

        # 5. General Exception
        mock_client.generate.side_effect = Exception("generic")
        with self.assertRaises(RuntimeError):
            provider.generate_quiz("text", 1, "MCQ", "Basic", "1-line")

    # Test quiz_generator local generators
    def test_quiz_generator_locals(self):
        from quiz_generator import (
            _generate_local_mcq,
            _generate_local_vsa,
            _generate_local_short_answer,
            _generate_local_long_answer,
            _generate_local_essay,
            generate_quiz,
        )

        s = "Python is a programming language. It was created by Guido van Rossum."
        self.assertIsNotNone(_generate_local_mcq(s, "Basic", "1-line", 1))
        self.assertIsNotNone(_generate_local_mcq(s, "Intermediate", "2-line", 1))
        self.assertIsNotNone(_generate_local_mcq(s, "Pro", "Detailed", 1))
        self.assertIsNotNone(_generate_local_mcq(s, "Pro", "Essay", 1))

        self.assertIsNotNone(_generate_local_vsa(s, "Basic", "1-line", 1))
        self.assertIsNotNone(_generate_local_vsa(s, "Intermediate", "2-line", 1))
        self.assertIsNotNone(_generate_local_vsa(s, "Pro", "Detailed", 1))

        self.assertIsNotNone(_generate_local_short_answer(s, "Basic", "1-line", 1))
        self.assertIsNotNone(
            _generate_local_short_answer(s, "Intermediate", "2-line", 1)
        )
        self.assertIsNotNone(_generate_local_short_answer(s, "Pro", "Detailed", 1))

        self.assertIsNotNone(_generate_local_long_answer(s, "Basic", "1-line", 1))
        self.assertIsNotNone(
            _generate_local_long_answer(s, "Intermediate", "2-line", 1)
        )
        self.assertIsNotNone(_generate_local_long_answer(s, "Pro", "Essay", 1))

        self.assertIsNotNone(_generate_local_essay(s, "Basic", "1-line", 1))
        self.assertIsNotNone(_generate_local_essay(s, "Intermediate", "2-line", 1))
        self.assertIsNotNone(_generate_local_essay(s, "Pro", "default", 1))

        # Empty sentences fallback
        self.assertIsNotNone(generate_quiz("Short text", 1))

        # Unsupported type warning path
        self.assertIsNotNone(
            generate_quiz(
                "This is a very long sentence. It has more than fifteen characters.",
                1,
                question_type="UNSUPPORTED",
            )
        )

    # Test generate_ai_quiz
    @patch("google.genai.Client")
    def test_generate_ai_quiz(self, mock_client_class):
        from quiz_generator import generate_ai_quiz

        mock_client = MagicMock()
        mock_response = MagicMock()
        mock_response.text = '{"quiz": [{"question": "q", "type": "MCQ", "answer": "a", "options": ["o1", "o2"], "correct_index": 0}]}'
        mock_client.models.generate_content.return_value = mock_response
        mock_client_class.return_value = mock_client

        res = generate_ai_quiz("text", 1, "MCQ", "Basic", "1-line", "api_key")
        self.assertEqual(len(res), 1)

    # Test generate_ai_quiz error paths
    @patch("google.genai.Client")
    def test_generate_ai_quiz_errors(self, mock_client_class):
        from quiz_generator import generate_ai_quiz
        import google.api_core.exceptions as google_exceptions

        mock_client = MagicMock()
        mock_client_class.return_value = mock_client

        # Empty api key
        with self.assertRaises(ValueError):
            generate_ai_quiz("text", 1, "MCQ", "Basic", "1-line", "")

        # ResourceExhausted exception
        mock_client.models.generate_content.side_effect = (
            google_exceptions.ResourceExhausted("rate limit")
        )
        with self.assertRaises(RuntimeError):
            generate_ai_quiz("text", 1, "MCQ", "Basic", "1-line", "api_key")

        # General exception
        mock_client.models.generate_content.side_effect = Exception("error")
        with self.assertRaises(RuntimeError):
            generate_ai_quiz("text", 1, "MCQ", "Basic", "1-line", "api_key")

    # Test generate_ollama_quiz wrapper
    @patch("quiz_generator.OllamaProvider")
    def test_generate_ollama_quiz(self, mock_ollama_prov):
        from quiz_generator import generate_ollama_quiz

        mock_prov = MagicMock()
        mock_prov.generate_quiz.return_value = [{"question": "q"}]
        mock_ollama_prov.return_value = mock_prov
        res = generate_ollama_quiz("text", 1, "MCQ", "Basic", "1-line")
        self.assertEqual(res, [{"question": "q"}])

    # Test document_parser external handlers
    @patch("document_parser.PdfReader")
    def test_document_parsers(self, mock_pdfreader):
        from document_parser import extract_text_from_file

        # 1. PDF standard
        mock_reader = MagicMock()
        mock_page = MagicMock()
        mock_page.extract_text.return_value = "pdf text"
        mock_reader.pages = [mock_page]
        mock_pdfreader.return_value = mock_reader
        res = extract_text_from_file(b"pdf data", "file.pdf")
        self.assertIn("pdf text", res)

        # 2. PDF OCR fallback
        mock_page.extract_text.return_value = ""
        mock_doc = MagicMock()
        mock_page_fitz = MagicMock()
        mock_pix = MagicMock()
        mock_pix.width = 100
        mock_pix.height = 100
        mock_pix.samples = b"\x00" * 30000
        mock_page_fitz.get_pixmap.return_value = mock_pix
        mock_doc.__len__.return_value = 1
        mock_doc.load_page.return_value = mock_page_fitz

        # Configure fitz mock
        import fitz

        fitz.open.return_value = mock_doc

        # Configure pytesseract mock
        import pytesseract

        pytesseract.image_to_string.return_value = "ocr text"

        res = extract_text_from_file(b"pdf data", "file.pdf")
        self.assertIn("ocr text", res)

        # 3. DOCX
        import docx

        mock_doc_obj = MagicMock()
        mock_para = MagicMock()
        mock_para.text = "docx text"
        mock_doc_obj.paragraphs = [mock_para]
        docx.Document.return_value = mock_doc_obj
        res = extract_text_from_file(b"docx data", "file.docx")
        self.assertEqual(res, "docx text")

        # Test docx import error
        real_docx = sys.modules.get("docx")
        sys.modules["docx"] = None
        try:
            res = extract_text_from_file(b"data", "file.docx")
            self.assertEqual(res, "[DOCX parser not installed]")
        finally:
            sys.modules["docx"] = real_docx

        # Test docx general exception
        docx.Document.side_effect = Exception("error")
        res = extract_text_from_file(b"docx data", "file.docx")
        self.assertEqual(res, "")
        docx.Document.side_effect = None

        # 4. PPTX
        import pptx

        mock_pres_obj = MagicMock()
        mock_slide = MagicMock()
        mock_shape = MagicMock()
        mock_shape.has_text_frame = True
        mock_para_pptx = MagicMock()
        mock_run = MagicMock()
        mock_run.text = "pptx text"
        mock_para_pptx.runs = [mock_run]
        mock_shape.text_frame.paragraphs = [mock_para_pptx]
        mock_slide.shapes = [mock_shape]
        mock_pres_obj.slides = [mock_slide]
        pptx.Presentation.return_value = mock_pres_obj
        res = extract_text_from_file(b"pptx data", "file.pptx")
        self.assertEqual(res, "pptx text")

        # Test pptx import error
        real_pptx = sys.modules.get("pptx")
        sys.modules["pptx"] = None
        try:
            res = extract_text_from_file(b"data", "file.pptx")
            self.assertEqual(res, "[PPTX parser not installed]")
        finally:
            sys.modules["pptx"] = real_pptx

        # Test pptx general exception
        pptx.Presentation.side_effect = Exception("error")
        res = extract_text_from_file(b"data", "file.pptx")
        self.assertEqual(res, "")
        pptx.Presentation.side_effect = None

        # 5. XLSX
        import openpyxl

        mock_wb = MagicMock()
        mock_sheet = MagicMock()
        mock_cell = MagicMock()
        mock_cell.value = "xlsx text"
        mock_sheet.iter_rows.return_value = [[mock_cell]]
        mock_wb.worksheets = [mock_sheet]
        openpyxl.load_workbook.return_value = mock_wb
        res = extract_text_from_file(b"xlsx data", "file.xlsx")
        self.assertEqual(res, "xlsx text")

        # Test xlsx import error
        real_openpyxl = sys.modules.get("openpyxl")
        sys.modules["openpyxl"] = None
        try:
            res = extract_text_from_file(b"data", "file.xlsx")
            self.assertEqual(res, "[XLSX parser not installed]")
        finally:
            sys.modules["openpyxl"] = real_openpyxl

        # Test xlsx general exception
        openpyxl.load_workbook.side_effect = Exception("error")
        res = extract_text_from_file(b"data", "file.xlsx")
        self.assertEqual(res, "")
        openpyxl.load_workbook.side_effect = None

        # 6. Image (PNG)
        res = extract_text_from_file(b"image data", "file.png")
        self.assertEqual(res, "ocr text")

        # Test image import error
        real_pytess = sys.modules.get("pytesseract")
        sys.modules["pytesseract"] = None
        try:
            res = extract_text_from_file(b"data", "file.png")
            self.assertEqual(res, "[OCR not configured]")
        finally:
            sys.modules["pytesseract"] = real_pytess

        # Test image general exception
        pytesseract.image_to_string.side_effect = Exception("error")
        res = extract_text_from_file(b"data", "file.png")
        self.assertEqual(res, "")
        pytesseract.image_to_string.side_effect = None
