import io
import csv
import logging
from pypdf import PdfReader

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    "pdf",
    "txt",
    "csv",
    "md",
    "json",
    "py",
    "js",
    "html",
    "css",
    "png",
    "jpg",
    "jpeg",  # For OCR
    "docx",
    "pptx",
    "xlsx",  # For document intelligence
}

# Mapping from app language codes to Tesseract language codes
TESSERACT_LANG_MAP = {
    "en": "eng",
    "te": "tel",
    "hi": "hin",
    "kn": "kan",
    "ta": "tam",
}


def allowed_file(filename: str) -> bool:
    """Checks if the file extension is supported."""
    return (
        "." in filename and filename.rsplit(".", 1)[1].lower() in SUPPORTED_EXTENSIONS
    )


def extract_text_from_txt(file_content: bytes) -> str:
    """Tries to decode text content with common encodings."""
    for enc in ("utf-8", "latin-1", "cp1252"):
        try:
            return file_content.decode(enc)
        except UnicodeDecodeError:
            continue
    return file_content.decode("utf-8", errors="ignore")


def extract_text_from_csv(file_content: bytes) -> str:
    """Extracts content from a CSV file, preserving rows."""
    decoded = extract_text_from_txt(file_content)
    try:
        reader = csv.reader(decoded.strip().splitlines())
        return "\n".join([" ".join(row) for row in reader])
    except Exception:
        return decoded  # Fallback to plain text extraction


def extract_text_from_pdf(file_content: bytes, lang: str = "en") -> str:
    """Extracts text from a PDF file, with an OCR fallback for scanned documents."""
    text = ""
    try:
        reader = PdfReader(io.BytesIO(file_content))
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    except Exception as e:
        logger.error(f"PyPDF2 extraction failed: {e}")

    # If no text is extracted, it might be a scanned PDF. Try OCR.
    if not text.strip():
        logger.info("No text found with standard PDF extraction, attempting OCR...")
        try:
            # This library can handle the conversion and OCR internally if needed
            # For this case, let's stick to a more direct OCR approach if pymupdf fails
            # We will convert pages to images and then OCR them.
            import fitz  # PyMuPDF
            import pytesseract
            from PIL import Image

            doc = fitz.open(stream=file_content, filetype="pdf")
            ocr_text = []
            tess_lang = TESSERACT_LANG_MAP.get(lang, "eng")

            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap()
                img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
                ocr_text.append(pytesseract.image_to_string(img, lang=tess_lang))

            text = "\n".join(ocr_text)
            logger.info(f"Extracted {len(text)} characters via OCR.")

        except ImportError:
            logger.warning(
                "OCR dependencies not found. Install `PyMuPDF`, `pytesseract`, and `Tesseract` for scanned PDF support."
            )
            return "[Scanned PDF - OCR not configured]"
        except Exception as e:
            logger.error(f"PDF OCR process failed: {e}")
            return ""

    return text


def extract_text_from_image(file_content: bytes, lang: str = "en") -> str:
    """Extracts text from an image using OCR (requires tesseract)."""
    try:
        import pytesseract
        from PIL import Image

        tess_lang = TESSERACT_LANG_MAP.get(lang, "eng")
        image = Image.open(io.BytesIO(file_content))
        return pytesseract.image_to_string(image, lang=tess_lang)
    except ImportError:
        logger.warning("OCR requires `pytesseract` and a Tesseract installation.")
        return "[OCR not configured]"
    except Exception as e:
        logger.error(f"Image OCR failed: {e}")
        return ""


def extract_text_from_docx(file_content: bytes) -> str:
    """Extracts text from a DOCX file."""
    try:
        import docx

        doc = docx.Document(io.BytesIO(file_content))
        return "\n".join([para.text for para in doc.paragraphs])
    except ImportError:
        logger.warning("DOCX parsing requires the `python-docx` library.")
        return "[DOCX parser not installed]"
    except Exception as e:
        logger.error(f"DOCX extraction failed: {e}")
        return ""


def extract_text_from_pptx(file_content: bytes) -> str:
    """Extracts text from a PPTX file."""
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(file_content))
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        return "\n".join(text_runs)
    except ImportError:
        logger.warning("PPTX parsing requires the `python-pptx` library.")
        return "[PPTX parser not installed]"
    except Exception as e:
        logger.error(f"PPTX extraction failed: {e}")
        return ""


def extract_text_from_xlsx(file_content: bytes) -> str:
    """Extracts text from an XLSX file."""
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(filename=io.BytesIO(file_content), read_only=True)
        text_content = []
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows():
                row_text = [str(cell.value) for cell in row if cell.value is not None]
                if row_text:
                    text_content.append(" ".join(row_text))
        return "\n".join(text_content)
    except ImportError:
        logger.warning("XLSX parsing requires the `openpyxl` library.")
        return "[XLSX parser not installed]"
    except Exception as e:
        logger.error(f"XLSX extraction failed: {e}")
        return ""


def extract_text_from_file(file_content: bytes, filename: str, lang: str = "en") -> str:
    """
    Unified function to extract text from various file types.
    This acts as the content extraction pipeline.
    """
    ext = filename.rsplit(".", 1)[1].lower()

    if not allowed_file(filename):
        raise ValueError(f"Unsupported file type: .{ext}")

    # Handlers for language-sensitive and other complex extractions
    if ext == "pdf":
        return extract_text_from_pdf(file_content, lang=lang)
    if ext in ("png", "jpg", "jpeg"):
        return extract_text_from_image(file_content, lang=lang)

    # Standard handlers
    extraction_map = {
        "csv": extract_text_from_csv,
        "docx": extract_text_from_docx,
        "pptx": extract_text_from_pptx,
        "xlsx": extract_text_from_xlsx,
    }

    # Use specific extractor if available, otherwise fallback to generic text
    extractor = extraction_map.get(ext, extract_text_from_txt)

    return extractor(file_content)
